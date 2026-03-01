import os
import asyncio
from pptx import Presentation
from app.graph.ppt_generation_agent.tools.pptx_parser import extract_text_from_pptx
from app.graph.ppt_generation_agent.state import PPTGenerationState

async def test_flow():
    # 1. Create a dummy PPTX
    pptx_path = "test_input.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Existing Pitch Deck"
    subtitle = slide.placeholders[1]
    subtitle.text = "This is a very boring and non-premium presentation."
    prs.save(pptx_path)
    
    print(f"Created {pptx_path}")
    
    # 2. Extract text
    text = extract_text_from_pptx(pptx_path)
    print(f"Extracted text:\n{text}")
    
    # 3. Test State structure
    state: PPTGenerationState = {
        "research_data": text,
        "logo_path": None,
        "color_palette": None,
        "use_default_colors": True,
        "draft": None,
        "critique": None,
        "iteration": 0,
        "ppt_path": None,
        "mode": "edit"
    }
    
    print("State created successfully.")
    print(f"Mode set to: {state['mode']}")
    
    # Cleanup
    if os.path.exists(pptx_path):
        os.remove(pptx_path)
    print("Cleanup done.")

if __name__ == "__main__":
    asyncio.run(test_flow())
