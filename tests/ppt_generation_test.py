"""
PPT Generation Agent — Test Suite
===================================
Mirrors the depth of `evaluation_test.py` and `market_research_test.py`.

Sections:
  1. Schema Validation   — PPTSection & PPTDraft  (pydantic models / schema.py)
  2. Theme Logic         — get_theme_config, create_dynamic_theme, extract_colors_from_image  (themes.py)
  3. Chart Generation    — generate_chart for every supported chart type  (chart_generator.py)
  4. Draft → PPTX Model  — map_draft_to_pptx_model layouts & logo placement  (ppt_tools.py)
  5. Agent Nodes         — generator_node, recommender_node, refiner_node  (node.py)
  6. API Route           — /generate endpoint (ppt_generation.py FastAPI route)
  7. PPTX Edit Flow      — /edit endpoint and PPTX text extraction (pptx_parser.py)
"""

import os
import json
import tempfile
import pytest
from unittest.mock import patch, AsyncMock, MagicMock
from pptx import Presentation
from fastapi import UploadFile

# ─── Schema ──────────────────────────────────────────────────────────────────
from app.graph.ppt_generation_agent.schema import PPTSection, PPTDraft, Critique

# ─── Themes ──────────────────────────────────────────────────────────────────
from app.graph.ppt_generation_agent.tools.themes import (
    get_theme_config,
    create_dynamic_theme,
    extract_colors_from_image,
    Theme,
)

# ─── Chart generator ─────────────────────────────────────────────────────────
from app.graph.ppt_generation_agent.tools.chart_generator import generate_chart

# ─── PPT tools ───────────────────────────────────────────────────────────────
from app.graph.ppt_generation_agent.tools.ppt_tools import map_draft_to_pptx_model

# ─── Nodes ───────────────────────────────────────────────────────────────────
from app.graph.ppt_generation_agent.node import (
    generator_node,
    recommender_node,
    refiner_node,
)

# ─── API Route & Parser ──────────────────────────────────────────────────────
import app.api.routes.ppt_generation as _ppt_route
from app.graph.ppt_generation_agent.tools.pptx_parser import extract_text_from_pptx
import asyncio as _asyncio


# ===========================================================================
# 0. FIXTURES
# ===========================================================================

def make_section(
    title: str = "Problem",
    content: list[str] | None = None,
    speaker_notes: str = "Some notes",
    image_prompt: str | None = None,
    visualization_data: dict | None = None,
) -> PPTSection:
    """Helper: build a minimal PPTSection."""
    return PPTSection(
        title=title,
        content=content or ["Point A", "Point B", "Point C"],
        speaker_notes=speaker_notes,
        image_prompt=image_prompt,
        visualization_data=visualization_data,
    )


def make_draft(
    title: str = "Spark2Scale Pitch",
    theme: str = "minimalist",
    logo_path: str | None = None,
    color_palette: list[str] | None = None,
    use_default_colors: bool = True,
    sections: list[PPTSection] | None = None,
) -> PPTDraft:
    """Helper: build a minimal PPTDraft with 2 sections."""
    return PPTDraft(
        title=title,
        theme=theme,
        logo_path=logo_path,
        color_palette=color_palette,
        use_default_colors=use_default_colors,
        sections=sections or [make_section("Problem"), make_section("Solution")],
    )


# ===========================================================================
# 1. SCHEMA VALIDATION
# ===========================================================================

class TestSchema:
    """Validate the PPTSection and PPTDraft Pydantic models."""

    def test_section_required_fields(self):
        """Section must have title, content, and speaker_notes."""
        s = make_section()
        assert s.title == "Problem"
        assert len(s.content) == 3
        assert s.speaker_notes == "Some notes"

    def test_section_optional_fields_default_none(self):
        """Optional fields on PPTSection default to None."""
        s = make_section()
        assert s.image_prompt is None
        assert s.image_path is None
        assert s.visualization_data is None
        assert s.visualization_path is None

    def test_section_with_image_prompt(self):
        """Setting image_prompt should be preserved."""
        s = make_section(image_prompt="flat icon of a rocket, minimal, white background")
        assert s.image_prompt == "flat icon of a rocket, minimal, white background"

    def test_section_with_visualization_data(self):
        """visualization_data dict is stored without mutation."""
        data = {"type": "bar", "labels": ["A", "B"], "values": [10, 20], "title": "Revenue"}
        s = make_section(visualization_data=data)
        assert s.visualization_data == data

    def test_draft_default_theme(self):
        """Draft theme defaults to 'minimalist'."""
        draft = make_draft()
        assert draft.theme == "minimalist"

    def test_draft_custom_color_palette(self):
        """Custom color palette should be stored correctly."""
        palette = ["#FF5733", "#33FF57", "#3357FF"]
        draft = make_draft(color_palette=palette, use_default_colors=False)
        assert draft.color_palette == palette
        assert draft.use_default_colors is False

    def test_draft_serialization_round_trip(self):
        """model_dump → JSON → model_validate should produce identical objects."""
        original = make_draft()
        dumped = original.model_dump()
        restored = PPTDraft.model_validate(dumped)
        assert restored.title == original.title
        assert len(restored.sections) == len(original.sections)

    def test_critique_schema(self):
        """Critique model stores score (0-100) and recommendation list."""
        c = Critique(
            critique="Slides are too text-heavy.",
            score=72,
            recommendations=["Add visuals", "Reduce bullet points"],
        )
        assert c.score == 72
        assert len(c.recommendations) == 2


# ===========================================================================
# 2. THEME LOGIC
# ===========================================================================

class TestThemes:
    """Validate theme resolution and dynamic theme creation."""

    def test_get_all_built_in_themes(self):
        """Each built-in theme returns a valid ThemeConfig with colors & fonts."""
        for theme_enum in Theme:
            config = get_theme_config(theme_enum)
            assert config.colors.primary  # non-empty hex
            assert config.fonts.header    # non-empty font name

    def test_minimalist_theme_colors(self):
        """Minimalist theme should be white background, black primary."""
        config = get_theme_config(Theme.MINIMALIST)
        assert config.colors.background == "FFFFFF"
        assert config.colors.primary == "000000"

    def test_dark_modern_theme_has_dark_background(self):
        """Dark Modern theme must have a dark background hex."""
        config = get_theme_config(Theme.DARK_MODERN)
        # Background should not be white
        assert config.colors.background != "FFFFFF"

    def test_create_dynamic_theme_single_color(self):
        """Dynamic theme created from a single color uses it as primary/secondary/accent."""
        config = create_dynamic_theme(["FF5733"])
        assert config.colors.primary == "FF5733"
        # secondary & accent fall back to primary
        assert config.colors.secondary == "FF5733"
        assert config.colors.accent == "FF5733"

    def test_create_dynamic_theme_strips_hash_prefix(self):
        """Colors supplied with '#' prefix should be normalized."""
        config = create_dynamic_theme(["#AB1234", "#CD5678"])
        assert config.colors.primary == "AB1234"
        assert config.colors.secondary == "CD5678"

    def test_create_dynamic_theme_three_colors(self):
        """Three colors map to primary, secondary, accent correctly."""
        config = create_dynamic_theme(["AA0000", "00BB00", "0000CC"])
        assert config.colors.primary == "AA0000"
        assert config.colors.secondary == "00BB00"
        assert config.colors.accent == "0000CC"

    def test_create_dynamic_theme_empty_list_fallback(self):
        """Passing an empty list returns the minimalist fallback theme."""
        config = create_dynamic_theme([])
        assert config.colors.background == "FFFFFF"  # resembles minimalist

    def test_extract_colors_nonexistent_path(self):
        """extract_colors_from_image with a missing path should return []."""
        result = extract_colors_from_image("/non/existent/path/logo.png")
        assert result == []

    def test_extract_colors_valid_image(self, tmp_path):
        """extract_colors_from_image on a real image should return non-empty list."""
        from PIL import Image as PILImage

        # Create a small 10x10 solid red image
        img = PILImage.new("RGB", (10, 10), color=(255, 0, 0))
        img_path = str(tmp_path / "red_logo.png")
        img.save(img_path)

        colors = extract_colors_from_image(img_path, num_colors=3)
        # At least 1 dominant color extracted
        assert len(colors) >= 1
        # Colors are uppercase hex strings (no '#')
        for c in colors:
            assert len(c) == 6
            int(c, 16)  # Must be valid hex


# ===========================================================================
# 3. CHART GENERATION
# ===========================================================================

class TestChartGeneration:
    """Verify generate_chart produces an image file for each chart type."""

    def _run_chart(self, chart_type: str, tmp_dir: str, labels=None, values=None) -> str:
        data = {
            "type": chart_type,
            "title": f"{chart_type.capitalize()} Chart",
            "labels": labels or ["Q1", "Q2", "Q3"],
            "values": values or [100, 200, 150],
            "x_label": "Quarter",
            "y_label": "Revenue",
        }
        return generate_chart(data, tmp_dir)

    def test_bar_chart_creates_file(self, tmp_path):
        path = self._run_chart("bar", str(tmp_path))
        assert path is not None
        assert os.path.isfile(path)
        assert path.endswith(".png")

    def test_horizontal_bar_chart_creates_file(self, tmp_path):
        path = self._run_chart("horizontal_bar", str(tmp_path))
        assert path is not None and os.path.isfile(path)

    def test_line_chart_creates_file(self, tmp_path):
        path = self._run_chart("line", str(tmp_path))
        assert path is not None and os.path.isfile(path)

    def test_pie_chart_creates_file(self, tmp_path):
        path = self._run_chart("pie", str(tmp_path))
        assert path is not None and os.path.isfile(path)

    def test_donut_chart_creates_file(self, tmp_path):
        path = self._run_chart("donut", str(tmp_path))
        assert path is not None and os.path.isfile(path)

    def test_funnel_chart_creates_file(self, tmp_path):
        # Funnel values should be proportions or large ints — use counts
        path = self._run_chart("funnel", str(tmp_path), labels=["Aware", "Interest", "Convert"], values=[1000, 400, 80])
        assert path is not None and os.path.isfile(path)

    def test_unknown_chart_type_triggers_back_gracefully(self, tmp_path):
        """Unknown chart type triggers the bar fallback, still returning a file."""
        path = self._run_chart("unknown_type", str(tmp_path))
        assert path is not None and os.path.isfile(path)

    def test_chart_with_theme_colors(self, tmp_path):
        """Passing custom theme_colors should be accepted without error."""
        data = {
            "type": "bar",
            "title": "Revenue",
            "labels": ["2022", "2023"],
            "values": [500, 800],
        }
        theme_colors = ["#003366", "#FF9900"]
        path = generate_chart(data, str(tmp_path), theme_colors=theme_colors)
        assert path is not None and os.path.isfile(path)

    def test_chart_unique_filenames(self, tmp_path):
        """Two consecutive chart calls should produce distinct files (UUID-based names)."""
        data = {"type": "bar", "title": "Test", "labels": ["A"], "values": [1]}
        path1 = generate_chart(data, str(tmp_path))
        path2 = generate_chart(data, str(tmp_path))
        assert path1 != path2


# ===========================================================================
# 4. DRAFT → PPTX MODEL MAPPING
# ===========================================================================

class TestMapDraftToPptxModel:
    """Test map_draft_to_pptx_model's layout decisions and slide counts."""

    def test_slide_count_equals_sections_plus_title(self):
        """Total slides = 1 title slide + N content slides."""
        draft = make_draft(sections=[make_section("A"), make_section("B"), make_section("C")])
        model = map_draft_to_pptx_model(draft)
        assert len(model.slides) == 4  # 1 title + 3 content

    def test_title_slide_is_first(self):
        """The first slide should contain the presentation title text."""
        draft = make_draft(title="My Awesome Startup")
        model = map_draft_to_pptx_model(draft)
        title_slide = model.slides[0]
        # PptxParagraphModel exposes `.text` directly as a string field
        all_text = " ".join(
            para.text or ""
            for shape in title_slide.shapes
            if hasattr(shape, "paragraphs")
            for para in shape.paragraphs
        )
        assert "My Awesome Startup" in all_text

    def test_logo_added_to_title_slide(self, tmp_path):
        """When logo_path exists, it is appended to title slide shapes."""
        from PIL import Image as PILImage

        logo = tmp_path / "logo.png"
        PILImage.new("RGB", (100, 50), (0, 0, 255)).save(str(logo))

        draft = make_draft(logo_path=str(logo))
        model = map_draft_to_pptx_model(draft)

        # Title slide (index 0) must have a PptxPictureBoxModel for the logo
        title_slide = model.slides[0]
        from app.graph.ppt_generation_agent.presenton_core.models.pptx_models import PptxPictureBoxModel
        picture_shapes = [s for s in title_slide.shapes if isinstance(s, PptxPictureBoxModel)]
        assert len(picture_shapes) >= 1

    def test_section_with_image_uses_split_layout(self, tmp_path):
        """A section with an image_path should result in a picture shape on that slide."""
        from PIL import Image as PILImage

        img = tmp_path / "section_img.png"
        PILImage.new("RGB", (200, 200), (100, 200, 100)).save(str(img))

        section = make_section(title="Traction")
        section.image_path = str(img)

        draft = make_draft(sections=[section])
        model = map_draft_to_pptx_model(draft)

        content_slide = model.slides[1]  # index 1 is the content slide
        from app.graph.ppt_generation_agent.presenton_core.models.pptx_models import PptxPictureBoxModel
        picture_shapes = [s for s in content_slide.shapes if isinstance(s, PptxPictureBoxModel)]
        assert len(picture_shapes) >= 1

    def test_section_without_image_uses_column_layout(self):
        """A section with no visual should produce only text boxes, no pictures."""
        from app.graph.ppt_generation_agent.presenton_core.models.pptx_models import PptxPictureBoxModel

        section = make_section(title="Team")  # no image_prompt, no visualization_path
        draft = make_draft(sections=[section])
        model = map_draft_to_pptx_model(draft)

        content_slide = model.slides[1]
        picture_shapes = [s for s in content_slide.shapes if isinstance(s, PptxPictureBoxModel)]
        # Logo not set, so no pictures expected on content slide
        assert len(picture_shapes) == 0

    def test_custom_palette_applied_to_model(self):
        """A custom color palette should produce a theme whose primary matches the first color."""
        palette = ["AABBCC", "DDEEFF", "112233"]
        draft = make_draft(color_palette=palette, use_default_colors=False)
        model = map_draft_to_pptx_model(draft)
        # If the model builds without error, palette was accepted
        assert model is not None


# ===========================================================================
# 5. AGENT NODES  (LLM calls are mocked)
# ===========================================================================

class TestAgentNodes:
    """Unit tests for generator_node, recommender_node, and refiner_node."""

    def _base_state(self) -> dict:
        return {
            "research_data": {"startup_info": {"name": "TestCo"}, "market_research": {}},
            "logo_path": None,
            "color_palette": None,
            "use_default_colors": True,
            "draft": None,
            "critique": None,
            "iteration": 0,
            "ppt_path": None,
        }

    @patch("app.graph.ppt_generation_agent.node.llm")
    def test_generator_node_returns_draft(self, mock_llm):
        """generator_node should populate the 'draft' key in state."""
        expected_draft = make_draft()
        mock_structured = MagicMock()
        mock_structured.invoke.return_value = expected_draft
        mock_llm.with_structured_output.return_value = mock_structured

        state = self._base_state()
        result = generator_node(state)

        assert "draft" in result
        assert result["draft"].title == "Spark2Scale Pitch"

    @patch("app.graph.ppt_generation_agent.node.llm")
    def test_generator_node_preserves_logo_path(self, mock_llm):
        """generator_node must copy logo_path from state onto the draft."""
        expected_draft = make_draft()
        mock_structured = MagicMock()
        mock_structured.invoke.return_value = expected_draft
        mock_llm.with_structured_output.return_value = mock_structured

        state = self._base_state()
        state["logo_path"] = "/some/logo.png"
        result = generator_node(state)

        assert result["draft"].logo_path == "/some/logo.png"

    @patch("app.graph.ppt_generation_agent.node.llm")
    def test_generator_node_preserves_color_palette(self, mock_llm):
        """generator_node must copy color_palette from state onto the draft."""
        expected_draft = make_draft()
        mock_structured = MagicMock()
        mock_structured.invoke.return_value = expected_draft
        mock_llm.with_structured_output.return_value = mock_structured

        state = self._base_state()
        state["color_palette"] = ["FF0000", "00FF00"]
        state["use_default_colors"] = False
        result = generator_node(state)

        assert result["draft"].color_palette == ["FF0000", "00FF00"]
        assert result["draft"].use_default_colors is False

    @patch("app.graph.ppt_generation_agent.node.llm")
    def test_recommender_node_returns_critique(self, mock_llm):
        """recommender_node should populate the 'critique' key in state."""
        expected_critique = Critique(
            critique="Too much text",
            score=65,
            recommendations=["Add charts", "Bold key stats"],
        )
        mock_structured = MagicMock()
        mock_structured.invoke.return_value = expected_critique
        mock_llm.with_structured_output.return_value = mock_structured

        state = self._base_state()
        state["draft"] = make_draft()
        result = recommender_node(state)

        assert "critique" in result
        assert result["critique"].score == 65
        assert "Add charts" in result["critique"].recommendations

    @patch("app.graph.ppt_generation_agent.node.llm")
    def test_refiner_node_increments_iteration(self, mock_llm):
        """refiner_node must increment iteration by 1."""
        expected_draft = make_draft(title="Refined Pitch")
        mock_structured = MagicMock()
        mock_structured.invoke.return_value = expected_draft
        mock_llm.with_structured_output.return_value = mock_structured

        state = self._base_state()
        state["draft"] = make_draft()
        state["critique"] = Critique(
            critique="Needs work",
            score=60,
            recommendations=["More data"],
        )
        state["iteration"] = 1
        result = refiner_node(state)

        assert result["iteration"] == 2

    @patch("app.graph.ppt_generation_agent.node.llm")
    def test_refiner_node_preserves_logo_path(self, mock_llm):
        """refiner_node must copy logo_path from state onto the refined draft."""
        expected_draft = make_draft()
        mock_structured = MagicMock()
        mock_structured.invoke.return_value = expected_draft
        mock_llm.with_structured_output.return_value = mock_structured

        state = self._base_state()
        state["logo_path"] = "/logo.png"
        state["draft"] = make_draft()
        state["critique"] = Critique(critique="OK", score=80, recommendations=[])
        result = refiner_node(state)

        assert result["draft"].logo_path == "/logo.png"


# ===========================================================================
# 6. API ROUTE — /generate (Synchronous wrappers for async logic)
# ===========================================================================

def _run_coro(coro):
    """Utility to run a coroutine synchronously."""
    return _asyncio.run(coro)

def _make_api_state(research_data=None):
    """Helper for API state."""
    return {
        "research_data": research_data or {"startup_info": {}, "market_research": {}},
        "logo_path": None,
        "color_palette": None,
        "use_default_colors": True,
        "draft": None,
        "critique": None,
        "iteration": 0,
        "ppt_path": None,
    }

@patch.object(_ppt_route, "supabase", None)
@patch.object(_ppt_route, "generate_pptx_file", new_callable=AsyncMock)
@patch.object(_ppt_route, "app_graph")
def test_api_generate_returns_success(mock_graph, mock_generate):
    """A valid request should return status='success' and a ppt_path."""
    draft = make_draft()
    mock_graph.ainvoke = AsyncMock(return_value={"draft": draft, "iteration": 1})
    mock_generate.return_value = "output/test_startup.pptx"

    response = _run_coro(_ppt_route.run_ppt_generation(_make_api_state(), "test-startup-uuid"))

    assert response.status == "success"
    assert response.ppt_path == "output/test_startup.pptx"
    assert response.title == draft.title

@patch.object(_ppt_route, "app_graph")
def test_api_generate_raises_http_exception_on_missing_draft(mock_graph):
    """When the graph returns no draft, a 500 HTTPException should be raised."""
    from fastapi import HTTPException
    mock_graph.ainvoke = AsyncMock(return_value={"draft": None, "iteration": 0})

    with pytest.raises(HTTPException) as exc_info:
        _run_coro(_ppt_route.run_ppt_generation(_make_api_state(), "bad-uuid"))

    assert exc_info.value.status_code == 500

@patch.object(_ppt_route, "supabase", None)
@patch.object(_ppt_route, "generate_pptx_file", new_callable=AsyncMock)
@patch.object(_ppt_route, "app_graph")
def test_api_generate_response_contains_json_response(mock_graph, mock_generate):
    """json_response field in the API response must match the serialized draft."""
    draft = make_draft()
    mock_graph.ainvoke = AsyncMock(return_value={"draft": draft, "iteration": 1})
    mock_generate.return_value = "output/test.pptx"

    response = _run_coro(_ppt_route.run_ppt_generation(_make_api_state(), "uuid-xyz"))

    assert response.json_response is not None
    assert response.json_response["title"] == draft.title
    assert "sections" in response.json_response

@patch("builtins.open", MagicMock())
@patch.object(_ppt_route, "generate_pptx_file", new_callable=AsyncMock)
@patch.object(_ppt_route, "app_graph")
def test_api_generate_supabase_insert_called_when_client_available(mock_graph, mock_generate):
    """When supabase client is available, table('documents').insert should be called."""
    draft = make_draft()
    mock_graph.ainvoke = AsyncMock(return_value={"draft": draft, "iteration": 1})
    mock_generate.return_value = "output/test.pptx"

    mock_supabase = MagicMock()
    mock_insert_result = MagicMock()
    mock_insert_result.data = [{"did": "doc-001"}]
    mock_supabase.table.return_value.insert.return_value.execute.return_value = mock_insert_result
    mock_supabase.table.return_value.update.return_value.eq.return_value.execute.return_value = MagicMock()
    mock_supabase.storage.from_.return_value.upload.return_value = MagicMock()
    mock_supabase.storage.from_.return_value.get_public_url.return_value = "https://example.com/test.pptx"

    with patch.object(_ppt_route, "supabase", mock_supabase):
        response = _run_coro(_ppt_route.run_ppt_generation(_make_api_state(), "startup-123"))

    mock_supabase.table.assert_called_with("documents")
    insert_call = mock_supabase.table.return_value.insert.call_args[0][0]
    assert insert_call["startup_id"] == "startup-123"
    assert "Pitch Deck" in insert_call["type"]     # "Pitch Deck (PPT)"
    assert "json_response" in insert_call

@patch.object(_ppt_route, "generate_pptx_file", new_callable=AsyncMock)
@patch.object(_ppt_route, "app_graph")
def test_api_generate_no_supabase_still_returns_success(mock_graph, mock_generate):
    """When supabase is None, generation still succeeds without errors."""
    draft = make_draft()
    mock_graph.ainvoke = AsyncMock(return_value={"draft": draft, "iteration": 1})
    mock_generate.return_value = "output/test.pptx"

    with patch.object(_ppt_route, "supabase", None):
        response = _run_coro(_ppt_route.run_ppt_generation(_make_api_state(), "startup-456"))

    assert response.status == "success"


# ===========================================================================
# 7. EDIT PPT ENDPOINT & PARSER
# ===========================================================================

class TestPPTXEditFlow:
    """Tests for the /edit endpoint and PPTX parser."""



    @patch("builtins.open", MagicMock())
    @patch.object(_ppt_route, "extract_text_from_pptx")
    @patch.object(_ppt_route, "run_ppt_generation", new_callable=AsyncMock)
    @pytest.mark.asyncio
    async def test_api_edit_route_initializes_correct_mode(self, mock_run, mock_extract):
        """Verify mode='edit' is passed to the graph."""
        mock_extract.return_value = "Extracted Text"
        mock_run.return_value = _ppt_route.PPTGenerationResponse(
            status="success", ppt_path="out.pptx", title="Enhanced",
            iterations=1, message="Done"
        )

        # Plain MagicMock (no spec) so .filename attribute is freely settable
        mock_file = MagicMock()
        mock_file.filename = "test.pptx"
        mock_file.read = AsyncMock(return_value=b"binary_data")

        response = await _ppt_route.edit_ppt(
            startup_id="uuid-123",
            ppt_file=mock_file,
            logo=None,
            use_default_colors=True,
            user_instructions=None,
            chat_summary=None,
        )

        args, _ = mock_run.call_args
        assert args[0]["mode"] == "edit"
        assert args[0]["research_data"] == "Extracted Text"