"""
tools.py — Pure utility classes (no LangChain / LangGraph dependencies).

Contains:
  - DocumentParser   : routes file types and parses them into annotated text.
  - SecurityGuardrails : sanitises PII and caps query length.
"""

import json
import os
import re

import fitz  # PyMuPDF
from pptx import Presentation


# ---------------------------------------------------------------------------
# Document parsing
# ---------------------------------------------------------------------------

class DocumentParser:
    """Handles routing, parsing, and text normalisation with injected spatial references."""

    @staticmethod
    def route_and_parse(file_data: str) -> str:
        """
        Accept either:
          a) a raw JSON string payload sent directly from the frontend, or
          b) a local file path (PDF / PPTX / JSON).
        Returns a single block of spatially-annotated plain text.
        """
        # 1. Detect inline JSON payload
        if file_data.strip().startswith("{") or file_data.strip().startswith("["):
            return DocumentParser._parse_json_string(file_data)

        # 2. Treat as a local file path
        ext = os.path.splitext(file_data)[1].lower()
        if ext == ".pdf":
            return DocumentParser._parse_pdf(file_data)
        elif ext == ".pptx":
            return DocumentParser._parse_pptx(file_data)
        elif ext == ".json":
            return DocumentParser._parse_json(file_data)
        else:
            raise ValueError(f"Unsupported file type or invalid payload format: {ext!r}")

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _normalize(text: str) -> str:
        return " ".join(text.split())

    @staticmethod
    def _parse_pdf(file_path: str) -> str:
        doc = fitz.open(file_path)
        context_blocks: list[str] = []
        for page_num, page in enumerate(doc, 1):
            lines = page.get_text("text").split("\n")
            for line_num, line in enumerate(lines, 1):
                clean_line = DocumentParser._normalize(line)
                if clean_line:
                    context_blocks.append(f"[Page {page_num}, Line {line_num}] {clean_line}")
        return "\n".join(context_blocks)

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        context_blocks: list[str] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            line_num = 1
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for line in shape.text.split("\n"):
                        clean_line = DocumentParser._normalize(line)
                        if clean_line:
                            context_blocks.append(
                                f"[Slide {slide_num}, Line {line_num}] {clean_line}"
                            )
                            line_num += 1
        return "\n".join(context_blocks)

    @staticmethod
    def _parse_json(file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return DocumentParser._format_json_data(data)

    @staticmethod
    def _parse_json_string(json_str: str) -> str:
        try:
            data = json.loads(json_str)
            return DocumentParser._format_json_data(data)
        except json.JSONDecodeError:
            raise ValueError("Provided payload looked like JSON but failed to decode.")

    @staticmethod
    def _format_json_data(data: dict | list) -> str:
        """Consistently format JSON (from a file or an inline string payload)."""
        lines = json.dumps(data, indent=2).split("\n")
        return "\n".join(f"[JSON Line {i}] {line}" for i, line in enumerate(lines, 1))


# ---------------------------------------------------------------------------
# Security guardrails
# ---------------------------------------------------------------------------

class SecurityGuardrails:
    """Input sanitisation using only the Python Standard Library."""

    def __init__(self) -> None:
        # Contact information
        self.email_pattern = re.compile(
            r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}"
        )
        self.phone_pattern = re.compile(
            r"\b(?:\+?\d{1,3}[-.\\s]?)?\(?\d{3}\)?[-.\\s]?\d{3}[-.\\s]?\d{4}\b"
        )
        # Financial / identity
        self.ssn_pattern = re.compile(r"\b\d{3}-\d{2}-\d{4}\b")
        self.credit_card_pattern = re.compile(r"\b(?:\d[ -]*?){13,16}\b")
        # URLs
        self.url_pattern = re.compile(
            r"https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+"
        )

    def sanitize_text(self, text: str) -> str:
        """Detect and redact sensitive structured entities."""
        clean = self.url_pattern.sub("[URL_REDACTED]", text)
        clean = self.email_pattern.sub("[EMAIL_REDACTED]", clean)
        clean = self.phone_pattern.sub("[PHONE_REDACTED]", clean)
        clean = self.ssn_pattern.sub("[SSN_REDACTED]", clean)
        clean = self.credit_card_pattern.sub("[CC_REDACTED]", clean)
        return clean

    def sanitize_query(self, query: str, max_length: int = 300) -> str:
        """Cap the query length to prevent DoS via massive context windows."""
        return query[:max_length].strip()