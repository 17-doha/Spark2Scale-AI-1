from pptx import Presentation
import os
from app.core.logger import get_logger

logger = get_logger(__name__)

def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extracts structured text from a PPTX file.
    Returns a string formatted slide by slide.
    """
    logger.info(f"Extracting text from PPTX: {pptx_path}")
    if not os.path.exists(pptx_path):
        logger.error(f"File not found: {pptx_path}")
        return ""

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error(f"Failed to load PPTX: {e}")
        return ""

    text_runs = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        
        # Try to find title
        if slide.shapes.title:
            slide_content.append(f"Title: {slide.shapes.title.text.strip()}")
        
        # Extract text from all other shapes
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        
        if slide_content:
            text_runs.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_content))
    
    return "\n\n".join(text_runs)
